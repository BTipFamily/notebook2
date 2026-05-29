"""Document parsing and chunking pipeline."""

import os
import io
import logging
from pathlib import Path
from typing import List, Tuple

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 150


def extract_text(file_path: str, file_type: str) -> Tuple[str, int]:
    """Returns (full_text, page_count). file_type is lowercase extension without dot."""
    file_path = Path(file_path)

    if file_type == "pdf":
        return _extract_pdf(file_path)
    elif file_type == "docx":
        return _extract_docx(file_path)
    elif file_type == "pptx":
        return _extract_pptx(file_path)
    elif file_type in ("xlsx", "xls", "csv"):
        return _extract_spreadsheet(file_path, file_type)
    elif file_type == "txt":
        text = file_path.read_text(encoding="utf-8", errors="replace")
        return text, 1
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(path: Path) -> Tuple[str, int]:
    import pdfplumber
    pages_text = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            # Also extract tables
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        text += "\n" + " | ".join(str(cell or "") for cell in row)
            pages_text.append(text)
    return "\n\n--- Page Break ---\n\n".join(pages_text), len(pages_text)


def _extract_docx(path: Path) -> Tuple[str, int]:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            paragraphs.append(row_text)
    return "\n\n".join(paragraphs), max(1, len(paragraphs) // 20)


def _extract_pptx(path: Path) -> Tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_parts = [f"--- Slide {i+1} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        slides_text.append("\n".join(slide_parts))
    return "\n\n".join(slides_text), len(prs.slides)


def _extract_spreadsheet(path: Path, file_type: str) -> Tuple[str, int]:
    if file_type == "csv":
        import csv
        rows = []
        with open(path, encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" | ".join(row))
        return "\n".join(rows), 1
    else:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = [f"=== Sheet: {sheet_name} ==="]
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append(" | ".join(str(cell or "") for cell in row))
            sheets.append("\n".join(rows))
        return "\n\n".join(sheets), len(wb.sheetnames)


def chunk_text(text: str, document_id: int, document_name: str) -> List[dict]:
    """Split text into overlapping chunks with metadata."""
    chunks = []
    sentences = _split_sentences(text)

    current_chunk = []
    current_len = 0
    chunk_idx = 0

    for sentence in sentences:
        sentence_len = len(sentence)
        if current_len + sentence_len > CHUNK_SIZE and current_chunk:
            chunk_text = " ".join(current_chunk)
            chunks.append({
                "text": chunk_text,
                "document_id": document_id,
                "document_name": document_name,
                "chunk_index": chunk_idx,
                "char_start": text.find(chunk_text[:50]) if chunk_text else 0,
            })
            chunk_idx += 1
            # Keep overlap: last few sentences
            overlap_sentences = []
            overlap_len = 0
            for s in reversed(current_chunk):
                if overlap_len + len(s) > CHUNK_OVERLAP:
                    break
                overlap_sentences.insert(0, s)
                overlap_len += len(s)
            current_chunk = overlap_sentences
            current_len = overlap_len

        current_chunk.append(sentence)
        current_len += sentence_len

    if current_chunk:
        chunks.append({
            "text": " ".join(current_chunk),
            "document_id": document_id,
            "document_name": document_name,
            "chunk_index": chunk_idx,
            "char_start": 0,
        })

    return chunks


def _split_sentences(text: str) -> List[str]:
    """Simple sentence splitter that also handles paragraphs."""
    import re
    # Split on sentence boundaries and paragraph breaks
    parts = re.split(r'(?<=[.!?])\s+|\n\n+', text)
    return [p.strip() for p in parts if p.strip() and len(p.strip()) > 10]
